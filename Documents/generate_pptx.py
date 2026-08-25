from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(filename="Smart_Tution_Fee_Portal_Presentation.pptx"):
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Smart Tution Fee Portal"
    subtitle.text = "Parul Institute of Computer Applications\nSemester V Project\n2026-27\n\nTeam members:\n[Your Enrollment No & Name]\nCompany Details:\n[Your Company Details]"

    # Helper function to create content slides
    def add_content_slide(title_text, content_text):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = title_text
        content.text = content_text
        # adjusting font size
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)

    # Slide 2: INDEX
    add_content_slide("INDEX", 
        "1. Abstract\n"
        "2. Comparison of New System with Existing System\n"
        "3. Technology and HW, SW Requirement Specification\n"
        "4. Modules and its short description\n"
        "5. Users and their role description\n"
        "6. Limitations\n"
        "7. Future Enhancement\n"
        "8. References & Bibliography"
    )

    # Slide 3: Abstract
    add_content_slide("Abstract",
        "The Smart Tution Fee Portal is a web-based platform designed to simplify "
        "and automate the collection, management, and tracking of tuition fees "
        "for educational institutions or private coaching centers. It provides a "
        "centralized database for administrators to easily monitor payments, pending dues, "
        "and generate financial reports. For students and parents, the portal offers a "
        "convenient and secure way to pay fees online, view payment history, and receive "
        "timely notifications for upcoming dues. The system reduces manual paperwork, "
        "minimizes errors, and enhances transparency in financial transactions."
    )

    # Slide 4: Comparison
    add_content_slide("Comparison of New System with Existing System",
        "Existing System:\n"
        "- Manual record keeping in registers or basic excel sheets.\n"
        "- Prone to human errors during calculation.\n"
        "- Time-consuming receipt generation.\n"
        "- Difficult to track defaulters and pending dues.\n\n"
        "New System:\n"
        "- Automated digital records and secure cloud database.\n"
        "- Error-free automated calculations.\n"
        "- Instant digital receipts and online secure payments.\n"
        "- Automated reminders and instant defaulter list generation."
    )

    # Slide 5: Tech & HW/SW
    add_content_slide("Technology and HW, SW Requirement Specification",
        "Software Requirements:\n"
        "- Frontend: HTML5, CSS3, JavaScript (React/Bootstrap)\n"
        "- Backend: Python (Django/Flask) or Node.js\n"
        "- Database: MySQL or PostgreSQL\n"
        "- Server: Apache / Nginx / Node\n\n"
        "Hardware Requirements:\n"
        "- Processor: Intel Core i3 or above\n"
        "- RAM: 4 GB or higher\n"
        "- Storage: 50 GB free space\n"
        "- Internet Connection for online features"
    )

    # Slide 6: Modules
    add_content_slide("Modules and its short description",
        "- Admin Module: Manage students, teachers, courses, fee structures, and view financial reports.\n"
        "- Student/Parent Module: View fee details, pay online, download receipts, check payment history.\n"
        "- Payment Gateway Module: Secure processing of online transactions.\n"
        "- Notification Module: Send automated SMS/Email reminders for pending fees.\n"
        "- Reporting Module: Generate daily/monthly revenue and defaulters reports."
    )

    # Slide 7: Features
    add_content_slide("Features and its short description",
        "- Online Fee Payment: 24x7 payment facility for students.\n"
        "- Automated Reminders: System alerts for upcoming or overdue fees.\n"
        "- Dashboard Analytics: Graphical representation of revenue and dues for admin.\n"
        "- Receipt Generation: One-click digital fee receipts.\n"
        "- Secure Authentication: Role-based access control and encrypted data."
    )

    # Slide 8: Users & Roles
    add_content_slide("Users and their role description",
        "- Admin: Has full control over the system. Can add/remove students, configure fees, and view all reports.\n"
        "- Teacher (Optional): View assigned students, track their fee status for class participation.\n"
        "- Student/Parent: Can log in to view their specific fee dues, make payments, and download receipts.\n"
        "- Accountant: Manages day-to-day transactions and manual payments if any."
    )

    # Slide 9: Limitations
    add_content_slide("Limitations",
        "- Requires continuous internet connectivity for online payments and real-time updates.\n"
        "- Initial setup and training required for staff moving from manual systems.\n"
        "- Dependency on third-party payment gateways for transaction success.\n"
        "- May require hardware upgrades for older institutional computers."
    )

    # Slide 10: Future Enhancement
    add_content_slide("Future Enhancement",
        "- Integration with Mobile App (Android/iOS) for easier access on smartphones.\n"
        "- AI-based predictive analysis for fee collection trends.\n"
        "- Integration with biometric attendance systems to link attendance with fee dues.\n"
        "- Support for multiple regional languages for wider accessibility."
    )

    # Slide 11: References
    add_content_slide("References & Bibliography",
        "- Web Development tutorials (w3schools.com, MDN Web Docs)\n"
        "- Payment Gateway API Documentation (e.g., Razorpay, Stripe)\n"
        "- Database Management Systems by Ramez Elmasri & Shamkant B. Navathe.\n"
        "- \"Software Engineering: A Practitioner's Approach\" by Roger S. Pressman.\n"
        "- IEEE papers on Online Payment Systems and Educational ERPs."
    )

    # Slide 12: Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Thank you !!!"
    
    prs.save(filename)
    print(f"Presentation saved as {filename}")

if __name__ == '__main__':
    create_presentation()
